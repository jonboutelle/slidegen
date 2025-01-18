"""
Presentation generator that creates PowerPoint presentations from JSON input files
using predefined templates and formats.
"""
from pptx import Presentation
import json
from typing import Dict, Optional
from enum import Enum
import os

class SlideFormat(Enum):
    """Available slide format templates"""
    IMAGE_AND_TEXT_SIDE_BY_SIDE = "IMAGE_AND_TEXT_SIDE_BY_SIDE"
    CENTERED_TEXT_ONLY = "CENTERED_TEXT_ONLY"

class PresentationGenerator:
    def __init__(self, template_path: str):
        """Initialize with path to template PPTX file"""
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template file not found: {template_path}")
        self.template_path = template_path
        self.prs = Presentation(template_path)

    def generate_from_json(self, json_path: str, output_path: str) -> None:
        """Generate a presentation from a JSON input file"""
        with open(json_path, 'r') as f:
            data = json.load(f)
        
        # Clear any existing slides
        while len(self.prs.slides) > 0:
            rId = self.prs.slides._sldIdLst[0].rId  # pylint: disable=protected-access
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[0]  # pylint: disable=protected-access
        
        for slide_data in data['slides']:
            self._add_slide(slide_data)
        
        self.prs.save(output_path)

    def _add_slide(self, slide_data: Dict) -> None:
        """Add a single slide based on the provided data"""
        format_str = slide_data.get('format')
        try:
            slide_format = SlideFormat(format_str)
        except ValueError:
            raise ValueError(f"Invalid slide format: {format_str}")

        text = slide_data.get('text')
        image_path = slide_data.get('image')

        # Select layout based on format
        if slide_format == SlideFormat.IMAGE_AND_TEXT_SIDE_BY_SIDE:
            slide_layout = self.prs.slide_layouts[8]  # 9th layout (0-based index)
        elif slide_format == SlideFormat.CENTERED_TEXT_ONLY:
            slide_layout = self.prs.slide_layouts[1]  # 2nd layout (0-based index)
        else:
            raise ValueError(f"Unhandled slide format: {format_str}")

        # Create new slide
        slide = self.prs.slides.add_slide(slide_layout)

        # Debug: Print placeholder information
        print(f"\nPlaceholders for {slide_format}:")
        for shape in slide.placeholders:
            print(f"- Index: {shape.placeholder_format.idx}, Type: {shape.name}")

        # Add content based on format
        if slide_format == SlideFormat.CENTERED_TEXT_ONLY:
            if text:
                # Use title placeholder (index 0)
                title = slide.placeholders[0]
                title.text = text
        elif slide_format == SlideFormat.IMAGE_AND_TEXT_SIDE_BY_SIDE:
            if text:
                # Add text to title placeholder (index 0)
                title = slide.placeholders[0]
                title.text = text
            
            if image_path:
                # Add image to picture placeholder (index 1)
                try:
                    img_placeholder = slide.placeholders[1]
                    img_placeholder.insert_picture(image_path)
                except (KeyError, FileNotFoundError) as e:
                    print(f"Warning: Could not add image {image_path}: {str(e)}")

def main():
    """Example usage"""
    try:
        generator = PresentationGenerator("template.pptx")
        generator.generate_from_json("example.json", "output.pptx")
        print("Successfully generated presentation: output.pptx")
    except Exception as e:
        print(f"Error generating presentation: {str(e)}")

if __name__ == "__main__":
    main()
